"""
PPTX Generator Utility — Reference implementation for the generate_pptx skill.

This module provides helper functions for building PowerPoint presentations
using python-pptx. It is loaded as a reference by the agent but the actual
code is executed in the sandboxed execution environment.

Dependencies (pre-installed in sandbox):
    python-pptx >= 0.6
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json


# ── Default theme ─────────────────────────────────────────────────────────────

THEME = {
    "header_bg":    RGBColor(0x1F, 0x4E, 0x79),
    "header_text":  RGBColor(0xFF, 0xFF, 0xFF),
    "body_text":    RGBColor(0x1F, 0x1F, 0x1F),
    "accent":       RGBColor(0x2E, 0x75, 0xB6),
    "slide_width":  Inches(13.33),
    "slide_height": Inches(7.5),
    "title_size":   Pt(36),
    "heading_size": Pt(28),
    "body_size":    Pt(18),
}


# ── Slide builders ────────────────────────────────────────────────────────────

def build_presentation(title, subtitle, slides):
    """
    Build a Presentation object from structured data.

    Args:
        title (str): Presentation title shown on the cover slide.
        subtitle (str): Subtitle shown on the cover slide.
        slides (list[dict]): List of slide dicts with keys:
            - title (str)
            - bullets (list[str])
            - notes (str, optional)

    Returns:
        Presentation: A ready-to-save python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = THEME["slide_width"]
    prs.slide_height = THEME["slide_height"]

    _add_title_slide(prs, title, subtitle)

    for slide_data in slides:
        _add_content_slide(
            prs,
            title=slide_data.get("title", ""),
            bullets=slide_data.get("bullets", []),
            notes=slide_data.get("notes", ""),
        )

    return prs


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_text(slide.shapes.title, title,
              color=THEME["header_text"], size=THEME["title_size"], bold=True)
    ph = slide.placeholders[1]
    _set_text(ph, subtitle,
              color=RGBColor(0xBF, 0xD7, 0xED), size=Pt(20))


def _add_content_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_text(slide.shapes.title, title,
              color=THEME["header_text"], size=THEME["heading_size"], bold=True)

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = THEME["body_size"]
            run.font.color.rgb = THEME["body_text"]

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _set_text(shape, text, color=None, size=None, bold=False):
    tf = shape.text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            if color:
                run.font.color.rgb = color
            if size:
                run.font.size = size
            run.font.bold = bold


# ── Entry point for standalone testing ───────────────────────────────────────

def generate(title, subtitle, slides, output_path):
    """
    Generate a PPTX file and save it.

    Args:
        title (str): Cover title.
        subtitle (str): Cover subtitle.
        slides (list[dict]): Slide data (title, bullets, notes).
        output_path (str): Path to write the .pptx file.

    Returns:
        dict: Result summary with status, path, and slide count.
    """
    prs = build_presentation(title, subtitle, slides)
    prs.save(output_path)
    result = {
        "status": "success",
        "output_path": output_path,
        "slide_count": len(slides) + 1,
    }
    print(json.dumps(result))
    return result
