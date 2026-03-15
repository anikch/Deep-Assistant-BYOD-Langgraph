import io
import logging
import httpx
from typing import Optional

logger = logging.getLogger(__name__)


def extract_from_pdf(file_path: str) -> str:
    """Extract text from PDF, with OCR fallback for image-only PDFs."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)

        full_text = "\n".join(text_parts).strip()

        # If very little text extracted, try OCR
        if len(full_text) < 100:
            logger.info(f"PDF has little text ({len(full_text)} chars), trying OCR")
            full_text = extract_pdf_with_ocr(file_path) or full_text

        return full_text
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return ""


def extract_pdf_with_ocr(file_path: str) -> str:
    """OCR fallback for image-based PDFs."""
    try:
        from PIL import Image
        import pytesseract
        import fitz  # PyMuPDF - try if available

        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img)
            text_parts.append(text)
        return "\n".join(text_parts)
    except ImportError:
        # fitz not available, skip OCR
        return ""
    except Exception as e:
        logger.error(f"OCR failed: {e}")
        return ""


def extract_from_pptx(file_path: str) -> str:
    """Extract text from all slides in a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {slide_idx + 1}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return ""


def extract_from_image(file_path: str) -> str:
    """OCR text from image file."""
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting image: {e}")
        return ""


def extract_from_txt(file_path: str) -> str:
    """Read plain text file."""
    try:
        import chardet

        with open(file_path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        return raw.decode(encoding, errors="replace")
    except Exception as e:
        logger.error(f"Error reading TXT: {e}")
        return ""


def extract_from_url(url: str) -> str:
    """Fetch URL and extract main text content using BeautifulSoup."""
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (compatible; AssistantPoC/1.0; +https://github.com)"
            )
        }
        response = httpx.get(url, headers=headers, timeout=30, follow_redirects=True)
        response.raise_for_status()

        from bs4 import BeautifulSoup

        soup = BeautifulSoup(response.text, "lxml")

        # Remove script/style/nav/footer
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form"]):
            tag.decompose()

        # Try to get main content
        main = soup.find("main") or soup.find("article") or soup.find("div", {"id": "content"})
        if main:
            text = main.get_text(separator="\n", strip=True)
        else:
            text = soup.get_text(separator="\n", strip=True)

        # Clean up excessive whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"Error fetching URL {url}: {e}")
        return ""


def extract_text(source_type: str, file_path: Optional[str] = None, url: Optional[str] = None) -> str:
    """Dispatch extraction based on source_type."""
    if source_type == "pdf":
        return extract_from_pdf(file_path)
    elif source_type == "pptx":
        return extract_from_pptx(file_path)
    elif source_type == "image":
        return extract_from_image(file_path)
    elif source_type in ("txt", "text"):
        return extract_from_txt(file_path)
    elif source_type == "url":
        return extract_from_url(url)
    else:
        logger.warning(f"Unknown source_type: {source_type}")
        return ""
